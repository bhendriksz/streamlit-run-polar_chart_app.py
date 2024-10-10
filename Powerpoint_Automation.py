import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from collections import defaultdict
import os
import colorsys

# Helper function to convert RGB values
def RGB(r, g, b):
    return RGBColor(r, g, b)

# Function to convert HSV to RGB and return as RGBColor
def hsv_to_rgb(h, s, v):
    r, g, b = colorsys.hsv_to_rgb(h, s, v)
    return RGB(int(r * 255), int(g * 255), int(b * 255))

# Define a mapping from department (two-letter code) to colors
department_colors = {}

# Function to assign a unique color to each department using HSV color space
def assign_department_color(department):
    if department not in department_colors:
        hue = len(department_colors) / 12.0  # Adjust denominator to control the spread of hues
        department_colors[department] = hsv_to_rgb(hue, 0.8, 0.9)  # Set saturation and value to fixed levels
    return department_colors[department]

# Function to add bullets to a slide with grouping by department
def add_bullets_to_slide_with_grouping(slide, projects, rows, cols, cell_width_cm, cell_height_cm, bol_diameter_pt):
    bullet_count = defaultdict(int)
    department_shapes = defaultdict(list)

    # Convert cell width and height from cm to inches
    cell_width_in = cell_width_cm / 2.54  # 1 inch = 2.54 cm
    cell_height_in = cell_height_cm / 2.54

    # Convert bullet (boll) diameter from pt to inches
    bol_diameter_in = bol_diameter_pt / 72.0  # 72 points = 1 inch

    # Loop over each project
    for afkorting, spf, project_title, project_type, source_slide_idx in projects:
        col_letter = spf[0]  # First letter for column
        row_number = int(spf[1:])  # Number for the row
        col_number = ord(col_letter.upper()) - ord('A') + 1  # Convert column letter to index

        # Only place the bullet if it fits within the grid
        if 1 <= row_number <= rows and 1 <= col_number <= cols:
            bullet_count[(row_number, col_number)] += 1
            bullet_idx = bullet_count[(row_number, col_number)] - 1

            # Maximum of 5 bullets per row, and 2 rows per cell
            visible_bullet_idx = bullet_idx % 10  # 0-9 visible bullets
            row_offset = (visible_bullet_idx // 5) * bol_diameter_in  # Offset for second row of bullets
            col_offset = (visible_bullet_idx % 5) * bol_diameter_in  # Offset for bullet position

            # Calculate the position on the slide in inches
            xPos = Inches(1 + (col_number - 1) * cell_width_in + col_offset)
            yPos = Inches(1 + (row_number - 1) * cell_height_in + row_offset)

            # Add bullet for the project
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, xPos, yPos, Inches(bol_diameter_in), Inches(bol_diameter_in))
            shape.text = afkorting

            # Assign a unique color to the department based on the two-letter code
            department_color = assign_department_color(afkorting[:2])
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = department_color

            # Set text font style and size
            text_frame = shape.text_frame
            text_frame.text = afkorting
            text_frame.paragraphs[0].font.name = "Arial"
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.italic = True

            # Collect shapes for grouping by department
            department_shapes[afkorting[:2]].append(shape)

# Function to add and run the VBA macro in the PowerPoint presentation
def add_and_run_macro(ppt_path, rows, cols, cell_width, cell_height, bol_diameter):
    pythoncom.CoInitialize()

    try:
        ppt_app = win32.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True  # Make sure PowerPoint is visible for debugging purposes

        # Open the presentation
        ppt_path = os.path.abspath(ppt_path)
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"The PowerPoint file was not found at: {ppt_path}")

        presentation = ppt_app.Presentations.Open(ppt_path)

        # Initialize a dictionary to store department-wise slides
        department_slides = defaultdict(list)
        department_initiatives_ideas = defaultdict(list)
        department_tasks = defaultdict(list)

        # Loop through all slides and shapes to find tables
        for sld in presentation.Slides:
            for shp in sld.Shapes:
                if shp.HasTable:
                    tbl = shp.Table
                    st.write(f"Table found on slide {sld.SlideIndex} with {tbl.Rows.Count} rows and {tbl.Columns.Count} columns")

                    afkorting_col = None
                    spf_col = None
                    title_col = None
                    type_col = None

                    # Identify the columns "AFKORTING", "SPF", "TITEL", and "PROJECT / IDEA / TASK"
                    for col in range(1, tbl.Columns.Count + 1):
                        header_text = tbl.Cell(1, col).Shape.TextFrame.TextRange.Text.upper()
                        st.write(f"Header found: {header_text}")
                        if header_text == "AFKORTING":
                            afkorting_col = col
                        elif header_text == "SPF":
                            spf_col = col
                        elif header_text == "TITEL":
                            title_col = col
                        elif header_text == "PROJECT / IDEA / TASK":
                            type_col = col

                    # Ensure the table has more than just a header row
                    if afkorting_col and spf_col and type_col and tbl.Rows.Count > 1:
                        # Loop through each row to group projects by department (Afkorting)
                        st.write(f"Processing table with {tbl.Rows.Count} rows.")
                        for row in range(2, tbl.Rows.Count + 1):
                            try:
                                afkorting = tbl.Cell(row, afkorting_col).Shape.TextFrame.TextRange.Text
                                spf = tbl.Cell(row, spf_col).Shape.TextFrame.TextRange.Text
                                project_title = tbl.Cell(row, title_col).Shape.TextFrame.TextRange.Text if title_col else ""
                                project_type = tbl.Cell(row, type_col).Shape.TextFrame.TextRange.Text

                                st.write(f"Row {row}: afkorting={afkorting}, spf={spf}, project_title={project_title}, project_type={project_type}")

                                if len(afkorting) >= 2:
                                    department = afkorting[:2]  # First two letters define the department
                                    department_slides[department].append((afkorting, spf, project_title, project_type, sld.SlideIndex))

                                    # Categorize into initiatives/ideas or tasks
                                    if project_type in ['Initiative', 'Idea']:
                                        department_initiatives_ideas[department].append((afkorting, spf, project_title, project_type, sld.SlideIndex))
                                    elif project_type == 'Task':
                                        department_tasks[department].append((afkorting, spf, project_title, project_type, sld.SlideIndex))
                            except Exception as e:
                                st.write(f"Error processing row {row}: {e}")

        # Create slide with all projects, using different colors for each department
        new_slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)  # Blank layout
        title_shape = new_slide.Shapes.AddTextbox(1, 10, 10, 500, 50)
        title_shape.TextFrame.TextRange.Text = "All Projects by Department"
        title_shape.TextFrame.TextRange.Font.Size = 24

        # Add all department projects in one slide
        all_projects = [project for projects in department_slides.values() for project in projects]
        add_bullets_to_slide_with_grouping(new_slide, all_projects, rows, cols, cell_width, cell_height, bol_diameter)

        # Slide for initiatives and ideas
        new_slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)  # Blank layout
        title_shape = new_slide.Shapes.AddTextbox(1, 10, 10, 500, 50)
        title_shape.TextFrame.TextRange.Text = "Initiatives and Ideas by Department"
        title_shape.TextFrame.TextRange.Font.Size = 24

        all_initiatives_ideas = [idea for ideas in department_initiatives_ideas.values() for idea in ideas]
        add_bullets_to_slide_with_grouping(new_slide, all_initiatives_ideas, rows, cols, cell_width, cell_height, bol_diameter)

        # Slide for tasks
        new_slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)  # Blank layout
        title_shape = new_slide.Shapes.AddTextbox(1, 10, 10, 500, 50)
        title_shape.TextFrame.TextRange.Text = "Tasks by Department"
        title_shape.TextFrame.TextRange.Font.Size = 24

        all_tasks = [task for tasks in department_tasks.values() for task in tasks]
        add_bullets_to_slide_with_grouping(new_slide, all_tasks, rows, cols, cell_width, cell_height, bol_diameter)

        # Save the modified presentation
        output_path = os.path.join(os.path.dirname(ppt_path), "updated_presentation_departments.pptm")
        presentation.SaveAs(output_path)
        presentation.Close()
        ppt_app.Quit()

        return output_path

    except Exception as e:
        st.error(f"An error occurred: {e}")
        if ppt_app:
            ppt_app.Quit()
        raise e


# Streamlit UI setup
st.title("PowerPoint Macro Tool (Headless)")

uploaded_ppt = st.file_uploader("Upload your PowerPoint presentation", type=["pptx"])

if uploaded_ppt is not None:
    with open("uploaded_presentation.pptx", "wb") as f:
        f.write(uploaded_ppt.read())

    st.success("PowerPoint successfully uploaded!")

    rows = st.number_input("Number of rows", min_value=1, step=1, value=23)
    cols = st.number_input("Number of columns", min_value=1, step=1, value=15)
    cell_width_cm = st.number_input("Width of each cell (cm)", min_value=0.1, step=0.1, value=4.1)  # Set input in cm
    cell_height_cm = st.number_input("Height of each cell (cm)", min_value=0.1, step=0.1, value=1.78)  # Set input in cm
    bol_diameter_pt = st.number_input("Diameter of the bullet (pt)", min_value=1.0, step=0.5, value=9.0)  # Set input in pt

    if st.button("Run Macro"):
        output_ppt_path = add_and_run_macro("uploaded_presentation.pptx", rows, cols, cell_width_cm, cell_height_cm, bol_diameter_pt)
        st.success("Macro executed successfully!")
        with open(output_ppt_path, "rb") as f:
            st.download_button("Download Updated PowerPoint", f, "updated_presentation_departments.pptx", "application/vnd.ms-powerpoint")

